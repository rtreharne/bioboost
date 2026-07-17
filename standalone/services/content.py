import hashlib
import json
import re
from pathlib import Path
from typing import Any

from bs4 import BeautifulSoup
from django.conf import settings
from django.db import transaction
from markdown import markdown
from openai import OpenAI
from openpyxl import load_workbook
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from standalone.models import ContentAsset, ContentChunk, LearningObjective
from standalone.services.guidance import sanitize_assistant_guidance
from standalone.services.math_questions import extract_math_subtopic_objectives
from standalone.services.symbol_heuristics import derive_symbol_heuristics_for_objectives


SUPPORTED_EXTENSIONS = {
    ".html",
    ".docx",
    ".pdf",
    ".txt",
    ".r",
    ".py",
    ".ipynb",
    ".rmd",
    ".md",
    ".pptx",
    ".xlsx",
}

OBJECTIVE_ACTION_VERBS = {
    "analyse",
    "analyze",
    "apply",
    "argue",
    "assess",
    "calculate",
    "choose",
    "classify",
    "compare",
    "construct",
    "contrast",
    "create",
    "critique",
    "define",
    "demonstrate",
    "describe",
    "design",
    "develop",
    "discuss",
    "distinguish",
    "evaluate",
    "examine",
    "explain",
    "identify",
    "illustrate",
    "interpret",
    "investigate",
    "justify",
    "outline",
    "plan",
    "predict",
    "propose",
    "reflect",
    "select",
    "solve",
    "summarise",
    "summarize",
    "use",
}
OBJECTIVE_STOPWORDS = {
    "a",
    "an",
    "and",
    "as",
    "at",
    "by",
    "for",
    "from",
    "how",
    "in",
    "into",
    "of",
    "on",
    "or",
    "the",
    "their",
    "through",
    "to",
    "using",
    "with",
    "your",
}
OBJECTIVE_META_TERMS = {
    "chapter",
    "content",
    "course",
    "document",
    "handout",
    "lecture",
    "material",
    "notes",
    "overview",
    "section",
    "session",
    "slide",
    "slides",
    "topic",
    "unit",
    "week",
}
OBJECTIVE_GENERIC_NOUNS = {
    "area",
    "areas",
    "aspect",
    "aspects",
    "awareness",
    "concept",
    "concepts",
    "content",
    "document",
    "idea",
    "ideas",
    "issue",
    "issues",
    "material",
    "materials",
    "overview",
    "part",
    "parts",
    "principle",
    "principles",
    "theme",
    "themes",
    "topic",
    "topics",
    "understanding",
}
OBJECTIVE_WEAK_OPENINGS = {
    "appreciate",
    "be aware of",
    "explore",
    "know",
    "learn about",
    "understand",
}


class BlockContentGenerationError(RuntimeError):
    pass


def normalize_text(value: str) -> str:
    value = value.replace("\r\n", "\n").replace("\r", "\n")
    value = re.sub(r"[ \t]+", " ", value)
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()


def _plain_text_from_markdown(text: str) -> str:
    return BeautifulSoup(markdown(text), "html.parser").get_text("\n")


def _read_plain_text_file(file_path: Path) -> str:
    raw_bytes = file_path.read_bytes()
    for encoding in ("utf-8-sig", "utf-16", "utf-16-le", "utf-16-be", "cp1252", "latin-1"):
        try:
            decoded = raw_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
        normalized = normalize_text(decoded)
        if normalized:
            return normalized
    return normalize_text(raw_bytes.decode("utf-8", errors="ignore"))


def extract_text_from_asset(asset: ContentAsset) -> str:
    ext = asset.extension.lower()
    file_path = Path(asset.file.path)

    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported content type: {ext}")

    if ext == ".txt":
        return _read_plain_text_file(file_path)

    if ext in {".r", ".py"}:
        return normalize_text(file_path.read_text(encoding="utf-8", errors="ignore"))

    if ext in {".md", ".rmd"}:
        return normalize_text(_plain_text_from_markdown(file_path.read_text(encoding="utf-8", errors="ignore")))

    if ext == ".html":
        return normalize_text(BeautifulSoup(file_path.read_text(encoding="utf-8", errors="ignore"), "html.parser").get_text("\n"))

    if ext == ".ipynb":
        notebook = json.loads(file_path.read_text(encoding="utf-8", errors="ignore"))
        cells = []
        for cell in notebook.get("cells", []):
            source = "".join(cell.get("source", []))
            if cell.get("cell_type") == "markdown":
                cells.append(_plain_text_from_markdown(source))
            else:
                cells.append(source)
        return normalize_text("\n\n".join(cells))

    if ext == ".docx":
        document = Document(file_path)
        return normalize_text("\n".join(paragraph.text for paragraph in document.paragraphs))

    if ext == ".pdf":
        reader = PdfReader(str(file_path))
        page_count = len(reader.pages)
        if page_count <= 0:
            return ""
        from standalone.services.pdf_import import extract_pdf_page_range

        return extract_pdf_page_range(file_path, 1, page_count)

    if ext == ".pptx":
        deck = Presentation(str(file_path))
        slides = []
        for slide in deck.slides:
            slide_lines = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
            if slide_lines:
                slides.append("\n".join(slide_lines))
        return normalize_text("\n\n".join(slides))

    if ext == ".xlsx":
        workbook = load_workbook(str(file_path), data_only=True)
        sheets = []
        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                values = [str(cell).strip() for cell in row if cell not in (None, "")]
                if values:
                    rows.append(" | ".join(values))
            if rows:
                sheets.append(f"{sheet.title}\n" + "\n".join(rows))
        return normalize_text("\n\n".join(sheets))

    raise ValueError(f"Unsupported content type: {ext}")


def _split_long_segment(text: str, target_size: int) -> list[str]:
    if len(text) <= target_size:
        return [text]

    sentences = [segment.strip() for segment in re.split(r"(?<=[.!?])\s+", text) if segment.strip()]
    if len(sentences) > 1:
        chunks: list[str] = []
        current: list[str] = []
        current_len = 0
        for sentence in sentences:
            if len(sentence) > target_size:
                if current:
                    chunks.append(" ".join(current))
                    current = []
                    current_len = 0
                chunks.extend(_split_long_segment(sentence, target_size))
                continue
            projected_len = current_len + len(sentence) + (1 if current else 0)
            if current and projected_len > target_size:
                chunks.append(" ".join(current))
                current = [sentence]
                current_len = len(sentence)
            else:
                current.append(sentence)
                current_len = projected_len
        if current:
            chunks.append(" ".join(current))
        return chunks

    words = text.split()
    if not words:
        return []

    chunks = []
    current_words: list[str] = []
    current_len = 0
    for word in words:
        if len(word) > target_size:
            if current_words:
                chunks.append(" ".join(current_words))
                current_words = []
                current_len = 0
            chunks.extend([word[index : index + target_size] for index in range(0, len(word), target_size)])
            continue
        projected_len = current_len + len(word) + (1 if current_words else 0)
        if current_words and projected_len > target_size:
            chunks.append(" ".join(current_words))
            current_words = [word]
            current_len = len(word)
        else:
            current_words.append(word)
            current_len = projected_len
    if current_words:
        chunks.append(" ".join(current_words))
    return chunks


def chunk_text(text: str, target_size: int = 1200) -> list[str]:
    paragraphs = [segment.strip() for segment in text.split("\n\n") if segment.strip()]
    if not paragraphs:
        return []

    chunks = []
    current = []
    current_len = 0
    for paragraph in paragraphs:
        paragraph_segments = _split_long_segment(paragraph, target_size)
        for segment in paragraph_segments:
            projected_len = current_len + len(segment) + (2 if current else 0)
            if current and projected_len > target_size:
                chunks.append("\n\n".join(current))
                current = [segment]
                current_len = len(segment)
            else:
                current.append(segment)
                current_len = projected_len if current_len else len(segment)
    if current:
        chunks.append("\n\n".join(current))
    return chunks


def sanitize_learning_objective(text: str) -> str:
    value = text.strip().strip("\"'`")
    value = value.replace("\u2013", "-").replace("\u2014", "-").replace("\u2212", "-")
    value = value.replace("\u2018", "'").replace("\u2019", "'").replace("\u201c", '"').replace("\u201d", '"')
    value = re.sub(r"^\s*[•●◦▪·*-]+\s*", "", value)
    value = re.sub(r"^\s*(?:learning\s*objective|objective|outcome|lo)\s*[:#-]?\s*\d+[a-z]?\s*[:.)-]?\s*", "", value, flags=re.IGNORECASE)
    value = re.sub(r"^\s*\(?\d+[a-z]?\)?\s*[.)-]\s*", "", value, flags=re.IGNORECASE)
    value = re.sub(r"^\s*\(?[ivxlcdm]+\)?\s*[.)-]\s*", "", value, flags=re.IGNORECASE)
    value = re.sub(r"^\s*[A-Z]\s*[.)-]\s*", "", value)
    value = value.strip(" -:\t")
    value = re.sub(r"\s+", " ", value)
    value = re.sub(r"[;:,.\-]+$", "", value)
    if not value:
        return ""
    if value[0].isalpha():
        value = value[0].upper() + value[1:]
    return value


def sanitize_summary(text: str) -> str:
    value = normalize_text(text)
    value = value.replace("\u2013", "-").replace("\u2014", "-")
    value = value.replace("\u2018", "'").replace("\u2019", "'").replace("\u201c", '"').replace("\u201d", '"')
    prefix_patterns = (
        r"^(?:the\s+)?(?:teaching\s+content|uploaded\s+teaching\s+material|uploaded\s+material|course\s+material|course\s+content|content|material)\s+(?:provides|offers|covers|introduces|explores|discusses|focuses\s+on|is\s+about)\s+",
        r"^(?:this\s+)?(?:teaching\s+content|uploaded\s+teaching\s+material|uploaded\s+material|course\s+material|course\s+content|content|material)\s+(?:provides|offers|covers|introduces|explores|discusses|focuses\s+on|is\s+about)\s+",
    )
    for pattern in prefix_patterns:
        value = re.sub(pattern, "", value, flags=re.IGNORECASE)
    value = re.sub(r"^(?:an?|the)\s+overview\s+of\s+", "", value, flags=re.IGNORECASE)
    value = re.sub(r"\s+([,.;:!?])", r"\1", value)
    value = re.sub(r"[ \t]+", " ", value)
    value = value.strip()
    if value and value[0].isalpha():
        value = value[0].upper() + value[1:]
    return value


def _dedupe_texts(items: list[str]) -> list[str]:
    seen: set[str] = set()
    seen_semantic: set[str] = set()
    deduped: list[str] = []
    for item in items:
        key = re.sub(r"[^a-z0-9]+", " ", item.lower()).strip()
        if not key or key in seen:
            continue
        semantic_key = _objective_similarity_key(item)
        if semantic_key and semantic_key in seen_semantic:
            continue
        seen.add(key)
        if semantic_key:
            seen_semantic.add(semantic_key)
        deduped.append(item)
    return deduped


def _objective_similarity_key(text: str) -> str:
    tokens = re.findall(r"[a-z0-9]+", sanitize_learning_objective(text).lower())
    if tokens and tokens[0] in OBJECTIVE_ACTION_VERBS:
        tokens = tokens[1:]
    meaningful_tokens = [token for token in tokens if token not in OBJECTIVE_STOPWORDS]
    return " ".join((meaningful_tokens or tokens)[:12])


def _objective_candidate_score(text: str) -> int:
    tokens = re.findall(r"[a-z0-9]+", sanitize_learning_objective(text).lower())
    if not tokens:
        return -10

    score = 0
    if tokens[0] in OBJECTIVE_ACTION_VERBS:
        score += 4
    if 5 <= len(tokens) <= 18:
        score += 2
    elif len(tokens) < 4 or len(tokens) > 28:
        score -= 2
    if any(token in {"compare", "contrast", "evaluate", "interpret", "justify", "analyse", "analyze", "apply"} for token in tokens):
        score += 1
    if any(token in OBJECTIVE_META_TERMS for token in tokens):
        score -= 3
    return score


def _source_keyword_set(text: str) -> set[str]:
    return {
        token
        for token in re.findall(r"[a-z0-9]+", normalize_text(text).lower())
        if len(token) >= 4 and token not in OBJECTIVE_STOPWORDS
    }


def _objective_validation_errors(objectives: list[str], source_text: str, max_items: int) -> list[str]:
    if len(objectives) < 3:
        return ["Return at least 3 learning objectives when content exists."]

    errors: list[str] = []
    source_keywords = _source_keyword_set(source_text)
    similarity_keys: set[str] = set()
    for index, objective in enumerate(objectives, start=1):
        lowered = objective.lower()
        tokens = re.findall(r"[a-z0-9]+", lowered)
        meaningful_tokens = [token for token in tokens if token not in OBJECTIVE_STOPWORDS]
        overlap = len(source_keywords & set(meaningful_tokens))
        similarity_key = _objective_similarity_key(objective)

        if similarity_key in similarity_keys:
            errors.append(f"Objective {index} duplicates another objective.")
        elif similarity_key:
            similarity_keys.add(similarity_key)

        if len(objective) < 32:
            errors.append(f"Objective {index} is too short and underspecified.")
        if len(objective) > 190:
            errors.append(f"Objective {index} is too long and should be more concise.")
        if overlap < 2:
            errors.append(f"Objective {index} is not specific enough to the source material.")
        if any(lowered.startswith(prefix) for prefix in OBJECTIVE_WEAK_OPENINGS):
            errors.append(f"Objective {index} uses a weak opening verb.")
        if any(term in lowered for term in ("overview", "course content", "teaching content", "uploaded material", "document structure")):
            errors.append(f"Objective {index} reads like a document summary rather than a learning objective.")
        if len(meaningful_tokens) < 4:
            errors.append(f"Objective {index} is too broad.")
        generic_count = sum(1 for token in meaningful_tokens if token in OBJECTIVE_GENERIC_NOUNS)
        if generic_count >= 2 and overlap <= 3:
            errors.append(f"Objective {index} relies on generic umbrella wording.")

    if len(objectives) < max_items and max_items >= 8 and len(source_keywords) >= 18:
        errors.append(f"Use more of the available objective slots when the source supports it, up to {max_items}.")
    return _dedupe_texts(errors)


def _block_content_validation_errors(summary: str, objectives: list[str], source_text: str, max_items: int) -> list[str]:
    errors = _objective_validation_errors(objectives, source_text, max_items)
    lowered_summary = summary.lower()
    if not summary.strip():
        errors.append("Return a non-empty summary.")
    if any(phrase in lowered_summary for phrase in ("teaching content", "uploaded material", "content provides", "material discusses")):
        errors.append("Rewrite the summary as a direct block description, not a description of the material.")
    return _dedupe_texts(errors)


def _revision_feedback_block(errors: list[str]) -> str:
    return "\n".join(f"- {error}" for error in errors if str(error).strip())


def _sentence_candidates_for_objectives(text: str, max_items: int) -> list[str]:
    sections = chunk_text(text, target_size=900) or [text]
    per_section_limit = max(1, min(3, max_items))
    candidates: list[str] = []

    for section in sections:
        scored_sentences: list[tuple[int, str]] = []
        for sentence in re.split(r"(?<=[.!?])\s+", section):
            stripped = sanitize_learning_objective(sentence)
            if 30 <= len(stripped) <= 180:
                scored_sentences.append((_objective_candidate_score(stripped), stripped))
        scored_sentences.sort(key=lambda item: (item[0], len(item[1])), reverse=True)
        candidates.extend(candidate for _score, candidate in scored_sentences[:per_section_limit])

    return candidates[: max_items * 2]


def _merge_objective_candidates(*objective_groups: list[str], max_items: int) -> list[str]:
    merged: list[str] = []
    for group in objective_groups:
        merged.extend(group)
    return _dedupe_texts(merged)[:max_items]


def _guidance_lines(text: str) -> list[str]:
    return [line.strip() for line in sanitize_assistant_guidance(text).split("\n") if line.strip()]


def _merge_guidance_texts(*texts: str) -> str:
    seen: set[str] = set()
    merged: list[str] = []
    for text in texts:
        for line in _guidance_lines(text):
            key = line.lower()
            if key in seen:
                continue
            seen.add(key)
            merged.append(line)
    return "\n".join(merged)


def _existing_objective_direction(block) -> list[str]:
    directions: list[str] = []
    for objective in block.learning_objectives.order_by("position", "pk"):
        text = sanitize_learning_objective(objective.text)
        guidance = sanitize_assistant_guidance(objective.assistant_guidance)
        if guidance:
            directions.append(f"{text} | Guidance: {guidance}")
        elif text:
            directions.append(text)
    return directions


def _block_objective_generation_guidance(block) -> str:
    course_guidance = ""
    try:
        course_guidance = block.course.config.assistant_guidance
    except Exception:  # noqa: BLE001
        course_guidance = ""
    block_guidance = ""
    try:
        block_guidance = block.config.assistant_guidance
    except Exception:  # noqa: BLE001
        block_guidance = ""
    return _merge_guidance_texts(course_guidance, block_guidance)


def _objective_budget_for_text(text: str, minimum: int = 6, maximum: int = 16) -> int:
    sections = max(1, len(chunk_text(text, target_size=900)))
    paragraphs = len([paragraph for paragraph in normalize_text(text).split("\n\n") if paragraph.strip()])
    line_candidates = len([line for line in text.splitlines() if sanitize_learning_objective(line)])
    return max(minimum, min(maximum, max(sections * 3, line_candidates, max(0, paragraphs // 2))))


def derive_learning_objectives(text: str, max_items: int = 10) -> list[str]:
    math_objectives = extract_math_subtopic_objectives(text, max_items=max_items)
    if math_objectives:
        return math_objectives[:max_items]

    candidates = []
    for line in text.splitlines():
        stripped = sanitize_learning_objective(line)
        if 30 <= len(stripped) <= 180 and stripped not in candidates:
            candidates.append(stripped)
        if len(candidates) >= max_items:
            break

    if len(candidates) < max_items:
        candidates.extend(_sentence_candidates_for_objectives(text, max_items=max_items))

    return _dedupe_texts(candidates)[:max_items]


def derive_learning_objectives_with_coverage(text: str, max_items: int = 10) -> list[str]:
    math_objectives = extract_math_subtopic_objectives(text, max_items=max_items)
    if math_objectives:
        return math_objectives[:max_items]

    sections = chunk_text(text, target_size=1200)
    if len(sections) <= 1:
        return derive_learning_objectives(text, max_items=max_items)

    per_section_limit = max(2, min(6, max_items // max(1, len(sections)) + 2))
    candidates: list[str] = []

    for section in sections:
        candidates.extend(derive_learning_objectives(section, max_items=per_section_limit))

    candidates.extend(derive_learning_objectives(text, max_items=max_items))
    return _dedupe_texts(candidates)[:max_items]


def _fallback_summary(text: str, max_sentences: int = 2, max_length: int = 320) -> str:
    normalized = normalize_text(text)
    sections = chunk_text(normalized, target_size=1200)
    representative_sentences: list[str] = []

    if len(sections) > 1:
        section_indexes = [0]
        if max_sentences > 1:
            section_indexes.append(len(sections) - 1)
        if max_sentences > 2 and len(sections) > 2:
            section_indexes.insert(1, len(sections) // 2)

        seen_sentences: set[str] = set()
        for index in section_indexes:
            section_sentences = [segment.strip() for segment in re.split(r"(?<=[.!?])\s+", sections[index]) if segment.strip()]
            if not section_sentences:
                continue
            sentence = section_sentences[0]
            sentence_key = sentence.lower()
            if sentence_key in seen_sentences:
                continue
            seen_sentences.add(sentence_key)
            representative_sentences.append(sentence)
            if len(representative_sentences) >= max_sentences:
                break

    if representative_sentences:
        summary = " ".join(representative_sentences).strip()
    else:
        sentences = [segment.strip() for segment in re.split(r"(?<=[.!?])\s+", normalized) if segment.strip()]
        summary = " ".join(sentences[:max_sentences]).strip()
    if not summary:
        summary = normalized[:max_length]
    if len(summary) > max_length:
        summary = summary[: max_length - 1].rsplit(" ", 1)[0].rstrip(",;:-") + "."
    if summary and summary[-1] not in ".!?":
        summary += "."
    return sanitize_summary(summary)


def _parse_json_object(raw_output: str) -> dict[str, Any]:
    if not raw_output:
        raise ValueError("OpenAI returned an empty payload.")
    try:
        return json.loads(raw_output)
    except json.JSONDecodeError:
        pass

    fenced_match = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", raw_output, re.DOTALL)
    if fenced_match:
        return json.loads(fenced_match.group(1))

    object_match = re.search(r"\{.*\}", raw_output, re.DOTALL)
    if object_match:
        return json.loads(object_match.group(0))

    raise ValueError("OpenAI did not return parseable JSON.")


def _block_content_response_schema(max_items: int) -> dict[str, Any]:
    return {
        "type": "object",
        "properties": {
            "summary": {"type": "string"},
            "learning_objectives": {
                "type": "array",
                "items": {"type": "string"},
                "minItems": 3,
                "maxItems": max_items,
            },
        },
        "required": ["summary", "learning_objectives"],
        "additionalProperties": False,
    }


def _learning_objective_generation_prompt(
    *,
    max_items: int,
    content: str,
    guidance: str = "",
    existing_objectives: list[str] | None = None,
    reduction_mode: bool = False,
    section_summaries: str = "",
    objective_lines: str = "",
    validation_feedback: str = "",
) -> str:
    guidance_block = guidance.strip() or "None provided."
    existing_objective_lines = "\n".join(f"- {item}" for item in (existing_objectives or []) if str(item).strip()) or "- None"
    if reduction_mode:
        content_block = f"""
Section summaries:
{section_summaries or "- None"}

Candidate learning objectives:
{objective_lines or "- None"}
""".strip()
    else:
        content_block = f"""
Content:
{content}
""".strip()
    feedback_block = ""
    if validation_feedback.strip():
        feedback_block = f"""
Revise the last attempt to fix these problems:
{validation_feedback}
""".strip()
    return f"""
Summarise the teaching content below for an educator-facing course authoring workflow.

Return strict JSON with exactly these keys:
- summary: one or two sentences in plain English
- learning_objectives: an array of 3 to {max_items} concise learning objectives

Author guidance (highest priority for audience, framing, and tone):
{guidance_block}

Existing learning-objective direction to preserve or refine when useful:
{existing_objective_lines}

Rules:
- follow the author guidance ahead of the source wording when choosing the audience, framing, tone, and emphasis
- use only the supplied source material for topic coverage and factual grounding
- no numbering or bullets in the output
- write the summary as a direct course/block description, not a description of the uploaded content
- do not use phrases like "the teaching content", "the uploaded material", "the content provides", "the textbook covers", or "this material discusses"
- no ambiguous references such as "this", "that", "it", or "etc."
- each learning objective should start with a clear action verb
- each learning objective must be student-focused and describe what a learner will do, understand, apply, evaluate, or reflect on
- each learning objective should focus on one concrete concept, mechanism, decision, comparison, case, or workflow
- prefer plain student-facing language over policy, governance, compliance, corporate, or institutional jargon
- if the source material is policy-heavy or administrative, translate it into clear study behaviours, decisions, and responsibilities for students
- keep learning objectives focused on the learner's actions and judgement, not the organisation's processes or document structure
- when rewording is needed, preserve the real topic but strip out jargon and corporate phrasing
- order the learning objectives as a sensible progression of themes or difficulty, not by the order of the uploaded documents or notes
- prefer breadth of coverage before nuance: capture the distinct themes, mechanisms, cases, decisions, or workflows in the source before adding finer-grained variants
- do not produce near-duplicate objectives that only swap the opening verb or lightly rephrase the same idea
- prefer specific subject nouns and mechanisms over umbrella wording such as topics, issues, areas, aspects, overview, principles, or understanding
- do not use weak openings such as "Understand", "Explore", "Know", or "Be aware of"
- if the source supports it, use more of the available objective slots rather than collapsing several substantive topics into one generic umbrella objective
- only return fewer than {max_items} objectives when the source genuinely contains fewer distinct learnable outcomes
- make every learning objective fully formed, distinct from the others, and easy to interpret on its own
- remove odd characters and formatting noise
- keep the wording specific enough to make sense on its own

{feedback_block}

{content_block}
""".strip()


def _openai_block_content_payload(
    text: str,
    max_items: int,
    *,
    guidance: str = "",
    existing_objectives: list[str] | None = None,
    validation_feedback: str = "",
) -> dict[str, Any]:
    client = OpenAI(api_key=settings.OPENAI_API_KEY)
    prompt = _learning_objective_generation_prompt(
        max_items=max_items,
        content=text,
        guidance=guidance,
        existing_objectives=existing_objectives,
        validation_feedback=validation_feedback,
    )
    response = client.responses.create(
        model=settings.OPENAI_MODEL,
        input=[
            {"role": "system", "content": [{"type": "input_text", "text": "Return only valid JSON."}]},
            {"role": "user", "content": [{"type": "input_text", "text": prompt}]},
        ],
        text={
            "format": {
                "type": "json_schema",
                "name": "block_content_generation",
                "strict": True,
                "schema": _block_content_response_schema(max_items),
            }
        },
    )
    return _parse_json_object((getattr(response, "output_text", "") or "").strip())


def _sanitize_objective_candidates(items: list[Any], max_items: int) -> list[str]:
    if not isinstance(items, list):
        return []
    return _dedupe_texts(
        [
            sanitized
            for item in items
            if isinstance(item, str)
            for sanitized in [sanitize_learning_objective(item)]
            if sanitized
        ]
    )[:max_items]


def _openai_reduce_content_payload(
    section_summaries: list[str],
    objective_candidates: list[str],
    max_items: int,
    *,
    guidance: str = "",
    existing_objectives: list[str] | None = None,
    validation_feedback: str = "",
) -> dict[str, Any]:
    client = OpenAI(api_key=settings.OPENAI_API_KEY)
    summary_lines = "\n".join(f"- {summary}" for summary in section_summaries if summary.strip())
    objective_lines = "\n".join(f"- {objective}" for objective in objective_candidates if objective.strip())
    prompt = _learning_objective_generation_prompt(
        max_items=max_items,
        content="",
        guidance=guidance,
        existing_objectives=existing_objectives,
        reduction_mode=True,
        section_summaries=summary_lines or "- None",
        objective_lines=objective_lines or "- None",
        validation_feedback=validation_feedback,
    )
    response = client.responses.create(
        model=settings.OPENAI_MODEL,
        input=[
            {"role": "system", "content": [{"type": "input_text", "text": "Return only valid JSON."}]},
            {"role": "user", "content": [{"type": "input_text", "text": prompt}]},
        ],
        text={
            "format": {
                "type": "json_schema",
                "name": "block_content_generation",
                "strict": True,
                "schema": _block_content_response_schema(max_items),
            }
        },
    )
    return _parse_json_object((getattr(response, "output_text", "") or "").strip())


def summarize_block_content(
    text: str,
    max_items: int = 6,
    *,
    guidance: str = "",
    existing_objectives: list[str] | None = None,
    allow_fallback: bool = True,
) -> tuple[str, list[str]]:
    math_objectives = extract_math_subtopic_objectives(text, max_items=max_items)
    fallback_summary = _fallback_summary(text)
    fallback_objectives = math_objectives or derive_learning_objectives_with_coverage(text, max_items=max_items)

    summary = fallback_summary
    objectives = fallback_objectives

    if not settings.OPENAI_API_KEY:
        if allow_fallback:
            return summary, objectives
        raise BlockContentGenerationError(
            "AI block generation requires OPENAI_API_KEY. Set it or explicitly enable heuristic fallback."
        )

    try:
        sections = chunk_text(text, target_size=5000)
        validation_errors: list[str] = []
        if len(sections) <= 1:
            payload = {}
            for _attempt in range(3):
                payload = _openai_block_content_payload(
                    text,
                    max_items,
                    guidance=guidance,
                    existing_objectives=existing_objectives,
                    validation_feedback=_revision_feedback_block(validation_errors),
                )
                summary = sanitize_summary(str(payload.get("summary", "")))
                generated_objectives = _sanitize_objective_candidates(payload.get("learning_objectives", []), max_items)
                validation_errors = _block_content_validation_errors(summary, generated_objectives, text, max_items)
                if not validation_errors:
                    break
            if validation_errors:
                raise ValueError("; ".join(validation_errors))
        else:
            per_section_limit = max(4, min(8, max_items // max(1, len(sections)) + 2))
            section_summaries: list[str] = []
            section_objectives: list[str] = []
            for section in sections:
                section_payload = _openai_block_content_payload(
                    section,
                    per_section_limit,
                    guidance=guidance,
                    existing_objectives=existing_objectives,
                )
                section_summary = sanitize_summary(str(section_payload.get("summary", "")))
                if section_summary:
                    section_summaries.append(section_summary)
                section_objectives.extend(
                    _sanitize_objective_candidates(section_payload.get("learning_objectives", []), per_section_limit)
                )
            merged_candidates = _dedupe_texts(section_objectives)
            payload = {}
            for _attempt in range(3):
                payload = _openai_reduce_content_payload(
                    section_summaries,
                    merged_candidates,
                    max_items,
                    guidance=guidance,
                    existing_objectives=existing_objectives,
                    validation_feedback=_revision_feedback_block(validation_errors),
                )
                summary = sanitize_summary(str(payload.get("summary", "")))
                generated_objectives = _sanitize_objective_candidates(payload.get("learning_objectives", []), max_items)
                validation_errors = _block_content_validation_errors(summary, generated_objectives, text, max_items)
                if not validation_errors:
                    break
            if validation_errors:
                raise ValueError("; ".join(validation_errors))

        summary = sanitize_summary(str(payload.get("summary", ""))) or fallback_summary
        generated_objectives = _sanitize_objective_candidates(payload.get("learning_objectives", []), max_items)
        if not generated_objectives:
            raise ValueError("OpenAI did not return usable learning objectives.")
        objectives = generated_objectives
    except Exception as exc:  # noqa: BLE001
        if allow_fallback:
            return fallback_summary, fallback_objectives
        raise BlockContentGenerationError(f"AI block generation failed: {exc}") from exc

    return summary, objectives


def _asset_text_for_regeneration(asset: ContentAsset) -> str:
    return asset.extracted_text or extract_text_from_asset(asset)


def resequence_learning_objectives(block, objectives: list[LearningObjective] | None = None) -> list[LearningObjective]:
    objectives = objectives or list(block.learning_objectives.order_by("position", "pk"))
    for index, objective in enumerate(objectives, start=1):
        objective.position = index
        objective.code = f"{block.order}.{index}"
    if objectives:
        LearningObjective.objects.bulk_update(objectives, ["position", "code"])
    return objectives


@transaction.atomic
def delete_block_and_resequence(block) -> None:
    course = block.course

    for asset in block.assets.all():
        asset.file.delete(save=False)

    block.delete()

    remaining_blocks = list(course.blocks.order_by("order", "created_at", "pk"))
    for index, remaining_block in enumerate(remaining_blocks, start=1):
        remaining_block.order = index
    if remaining_blocks:
        type(remaining_blocks[0]).objects.bulk_update(remaining_blocks, ["order"])

    for remaining_block in remaining_blocks:
        resequence_learning_objectives(remaining_block)

    course_fragments = [remaining_block.summary for remaining_block in remaining_blocks if remaining_block.summary.strip()]
    if course_fragments:
        course_summary, _ = summarize_block_content("\n\n".join(course_fragments), max_items=4)
        course.summary = course_summary
    else:
        course.summary = ""
    course.save(update_fields=["summary", "updated_at"])


@transaction.atomic
def move_course_block(block, direction: str) -> bool:
    blocks = list(block.course.blocks.order_by("order", "created_at", "pk"))
    try:
        current_index = next(index for index, item in enumerate(blocks) if item.pk == block.pk)
    except StopIteration:
        return False

    if direction == "up" and current_index > 0:
        swap_index = current_index - 1
    elif direction == "down" and current_index < len(blocks) - 1:
        swap_index = current_index + 1
    else:
        return False

    blocks[current_index], blocks[swap_index] = blocks[swap_index], blocks[current_index]
    for index, reordered_block in enumerate(blocks, start=1):
        reordered_block.order = index
    type(blocks[0]).objects.bulk_update(blocks, ["order"])

    for reordered_block in blocks:
        resequence_learning_objectives(reordered_block)
    return True


def _replace_block_objectives(block, source_asset: ContentAsset, objective_texts: list[str], context_text: str = "") -> int:
    existing = list(block.learning_objectives.order_by("position", "pk"))
    objective_heuristics = derive_symbol_heuristics_for_objectives(objective_texts, context_text)
    created_or_updated = 0
    for index, text in enumerate(objective_texts, start=1):
        heuristics = objective_heuristics[index - 1] if index - 1 < len(objective_heuristics) else {}
        defaults = {
            "course": block.course,
            "block": block,
            "source_asset": source_asset,
            "position": index,
            "code": f"{block.order}.{index}",
            "text": text,
            "symbol_heuristics": heuristics,
        }
        if index <= len(existing):
            objective = existing[index - 1]
            changed = False
            for field, value in defaults.items():
                if getattr(objective, field) != value:
                    setattr(objective, field, value)
                    changed = True
            if changed:
                objective.save(update_fields=["course", "block", "source_asset", "position", "code", "text", "symbol_heuristics", "updated_at"])
            created_or_updated += 1
            continue

        LearningObjective.objects.create(**defaults)
        created_or_updated += 1

    for objective in existing[len(objective_texts) :]:
        objective.delete()

    return created_or_updated


def refresh_learning_objective_symbol_heuristics(objective: LearningObjective, context_text: str = "") -> LearningObjective:
    heuristics = derive_symbol_heuristics_for_objectives([objective.text], context_text or "")[0]
    if objective.symbol_heuristics != heuristics:
        objective.symbol_heuristics = heuristics
        objective.save(update_fields=["symbol_heuristics", "updated_at"])
    return objective


def refresh_course_summary_from_blocks(course) -> bool:
    course_fragments = [block.summary for block in course.blocks.all() if block.summary.strip()]
    if not course_fragments:
        return False

    course_summary, _ = summarize_block_content("\n\n".join(course_fragments), max_items=4)
    course.summary = course_summary
    course.save(update_fields=["summary", "updated_at"])
    return True


def _clear_block_descriptions_and_objectives(block) -> None:
    block.summary = ""
    block.save(update_fields=["summary", "updated_at"])
    block.learning_objectives.all().delete()


def regenerate_block_descriptions_and_objectives(block, progress_callback=None, *, allow_fallback: bool = True) -> dict[str, int]:
    if progress_callback:
        progress_callback(15)
    assets = [asset for asset in block.assets.all() if asset.include_in_generation]
    if not assets:
        if progress_callback:
            progress_callback(100)
        return {"blocks": 0, "objectives": 0}

    if progress_callback:
        progress_callback(35)
    texts = [text for asset in assets if (text := _asset_text_for_regeneration(asset))]
    combined_text = normalize_text("\n\n".join(texts))
    if not combined_text:
        if progress_callback:
            progress_callback(100)
        return {"blocks": 0, "objectives": 0}

    objective_budget = _objective_budget_for_text(combined_text)
    generation_guidance = _block_objective_generation_guidance(block)
    existing_objectives = _existing_objective_direction(block)
    if progress_callback:
        progress_callback(65)
    try:
        block_summary, objectives = summarize_block_content(
            combined_text,
            max_items=objective_budget,
            guidance=generation_guidance,
            existing_objectives=existing_objectives,
            allow_fallback=allow_fallback,
        )
    except Exception:
        _clear_block_descriptions_and_objectives(block)
        raise
    if progress_callback:
        progress_callback(85)
    block.summary = block_summary
    block.save(update_fields=["summary", "updated_at"])
    objective_count = _replace_block_objectives(block, assets[0], objectives, combined_text)
    refresh_course_summary_from_blocks(block.course)
    if progress_callback:
        progress_callback(100)
    return {"blocks": 1, "objectives": objective_count}


@transaction.atomic
def move_learning_objective(objective: LearningObjective, direction: str) -> bool:
    objectives = list(objective.block.learning_objectives.order_by("position", "pk"))
    try:
        current_index = next(index for index, item in enumerate(objectives) if item.pk == objective.pk)
    except StopIteration:
        return False

    if direction == "up" and current_index > 0:
        swap_index = current_index - 1
    elif direction == "down" and current_index < len(objectives) - 1:
        swap_index = current_index + 1
    else:
        return False

    objectives[current_index], objectives[swap_index] = objectives[swap_index], objectives[current_index]
    resequence_learning_objectives(objective.block, objectives)
    return True


@transaction.atomic
def delete_learning_objective_and_resequence(objective: LearningObjective) -> None:
    block = objective.block
    objective.delete()
    resequence_learning_objectives(block)


def regenerate_course_descriptions_and_objectives(course, *, allow_fallback: bool = True) -> dict[str, int]:
    block_count = 0
    objective_count = 0

    for block in course.blocks.prefetch_related("assets", "learning_objectives").all():
        refreshed = regenerate_block_descriptions_and_objectives(block, allow_fallback=allow_fallback)
        block_count += refreshed["blocks"]
        objective_count += refreshed["objectives"]

    return {"blocks": block_count, "objectives": objective_count}


def generate_embeddings(texts: list[str]) -> list[list[float]]:
    if not settings.OPENAI_API_KEY:
        return [[] for _ in texts]
    try:
        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        response = client.embeddings.create(model=settings.OPENAI_EMBEDDING_MODEL, input=texts)
        return [item.embedding for item in response.data]
    except Exception:  # noqa: BLE001
        return [[] for _ in texts]


def ingest_content_asset(asset: ContentAsset, *, generate_objectives: bool = True, allow_fallback: bool = True) -> None:
    extracted_text = extract_text_from_asset(asset)
    asset.extracted_text = extracted_text
    asset.processing_status = ContentAsset.ProcessingStatus.PROCESSED
    asset.processing_error = ""
    asset.save(update_fields=["extracted_text", "processing_status", "processing_error", "updated_at"])

    asset.chunks.all().delete()
    asset.learning_objectives.all().delete()

    chunks = chunk_text(extracted_text)
    embeddings = generate_embeddings(chunks) if chunks else []
    for index, chunk in enumerate(chunks, start=1):
        ContentChunk.objects.create(
            asset=asset,
            course=asset.block.course,
            block=asset.block,
            ordinal=index,
            text=chunk,
            token_count=max(1, len(chunk.split())),
            embedding_model=settings.OPENAI_EMBEDDING_MODEL if embeddings and embeddings[index - 1] else "",
            embedding_vector=embeddings[index - 1] if embeddings else [],
            checksum=hashlib.sha256(chunk.encode("utf-8")).hexdigest(),
        )

    if not generate_objectives:
        return

    objective_budget = _objective_budget_for_text(extracted_text, minimum=4, maximum=16)
    try:
        _, objectives = summarize_block_content(
            extracted_text,
            max_items=objective_budget,
            guidance=_block_objective_generation_guidance(asset.block),
            allow_fallback=allow_fallback,
        )
    except Exception:
        _clear_block_descriptions_and_objectives(asset.block)
        raise
    objectives = [text for text in objectives if text]
    objective_heuristics = derive_symbol_heuristics_for_objectives(objectives, extracted_text)
    for index, objective in enumerate(objectives, start=1):
        LearningObjective.objects.create(
            course=asset.block.course,
            block=asset.block,
            source_asset=asset,
            position=index,
            code=f"{asset.block.order}.{index}",
            text=objective,
            symbol_heuristics=objective_heuristics[index - 1] if index - 1 < len(objective_heuristics) else {},
        )
